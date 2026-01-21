#!/usr/bin/env python3
"""
ATP Parser - Extract curriculum content from CAPS ATP PowerPoint files.
Generates structured JSON for seeding DynamoDB tables.
"""

import os
import re
import json
import uuid
from pptx import Presentation
from pathlib import Path


def generate_id(prefix: str) -> str:
    """Generate a unique ID with prefix."""
    return f"{prefix}-{uuid.uuid4().hex[:8]}"


def parse_filename(filename: str) -> dict:
    """Extract subject and grade info from filename."""
    # Patterns to match different filename formats
    # e.g., "2023 Grade 10-12 Mathematics ATP Mediation.pptx"
    # e.g., "2023 Grade 7 - 9 Natural Sciences ATP mediation.pptx"
    
    result = {
        "year": 2023,
        "subject": "",
        "grades": [],
        "raw_filename": filename
    }
    
    # Extract year
    year_match = re.search(r'(\d{4})', filename)
    if year_match:
        result["year"] = int(year_match.group(1))
    
    # Extract grade(s)
    grade_patterns = [
        r'Grade\s*(\d+)\s*[-–]\s*(\d+)',  # Grade 10-12
        r'Grades?\s*(\d+)\s*[-–]\s*(\d+)',  # Grades 10-12
        r'Grade\s*(\d+)',  # Single grade
        r'Grades?\s*(\d+)',
    ]
    
    for pattern in grade_patterns:
        match = re.search(pattern, filename, re.IGNORECASE)
        if match:
            if len(match.groups()) == 2:
                start, end = int(match.group(1)), int(match.group(2))
                result["grades"] = list(range(start, end + 1))
            else:
                result["grades"] = [int(match.group(1))]
            break
    
    # Extract subject name - remove common patterns
    subject = filename
    patterns_to_remove = [
        r'\d{4}',  # Year
        r'Grade\s*\d+\s*[-–]?\s*\d*',
        r'Grades?\s*\d+\s*[-–]?\s*\d*',
        r'ATP\s*Mediat?ion',
        r'Mediasie',
        r'\.pptx$',
        r'^\s+|\s+$',
    ]
    
    for pattern in patterns_to_remove:
        subject = re.sub(pattern, '', subject, flags=re.IGNORECASE)
    
    # Clean up extra spaces and special chars
    subject = re.sub(r'\s+', ' ', subject).strip()
    subject = re.sub(r'^[-–\s]+|[-–\s]+$', '', subject)
    
    result["subject"] = subject if subject else "Unknown"
    
    return result


def extract_text_from_slide(slide) -> list[str]:
    """Extract all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return texts


def extract_tables_from_slide(slide) -> list[list[list[str]]]:
    """Extract all tables from a slide."""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace('\n', ' ').replace('  ', ' ')
                    row_data.append(cell_text)
                table_data.append(row_data)
            tables.append(table_data)
    return tables


def extract_topics_from_tables(tables: list) -> list[dict]:
    """Extract topics from table data."""
    topics = []
    
    for table in tables:
        if not table:
            continue
            
        # Check for content/topic columns
        header = table[0] if table else []
        header_lower = [h.lower() for h in header]
        
        # Look for tables with topic/content information
        topic_col = -1
        term_col = -1
        
        for i, h in enumerate(header_lower):
            if 'content' in h or 'topic' in h:
                topic_col = i
            if 'term' in h:
                term_col = i
        
        if topic_col >= 0:
            for row in table[1:]:  # Skip header
                if len(row) > topic_col:
                    topic_name = row[topic_col].strip()
                    if topic_name and len(topic_name) > 3:
                        term = 0
                        if term_col >= 0 and len(row) > term_col:
                            term_match = re.search(r'(\d+)', row[term_col])
                            if term_match:
                                term = int(term_match.group(1))
                        
                        topics.append({
                            "topicName": topic_name[:200],  # Limit length
                            "term": term,
                            "context": "",  # To be filled later
                        })
    
    return topics


def parse_pptx_file(filepath: str) -> dict:
    """Parse a single PPTX file and extract curriculum data."""
    filename = os.path.basename(filepath)
    file_info = parse_filename(filename)
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"Error reading {filename}: {e}")
        return None
    
    all_topics = []
    all_text = []
    
    for slide in prs.slides:
        # Extract text for context
        texts = extract_text_from_slide(slide)
        all_text.extend(texts)
        
        # Extract tables for structured data
        tables = extract_tables_from_slide(slide)
        topics = extract_topics_from_tables(tables)
        all_topics.extend(topics)
    
    # De-duplicate topics
    seen = set()
    unique_topics = []
    for topic in all_topics:
        key = topic["topicName"].lower()
        if key not in seen:
            seen.add(key)
            unique_topics.append(topic)
    
    return {
        "filename": filename,
        "subject": file_info["subject"],
        "year": file_info["year"],
        "grades": file_info["grades"],
        "topics": unique_topics,
        "slideCount": len(prs.slides),
    }


def process_all_atps(downloads_dir: str) -> list[dict]:
    """Process all ATP PPTX files in the Downloads directory."""
    results = []
    
    pptx_files = list(Path(downloads_dir).glob("*.pptx"))
    print(f"Found {len(pptx_files)} PPTX files")
    
    for pptx_file in pptx_files:
        print(f"Processing: {pptx_file.name}")
        data = parse_pptx_file(str(pptx_file))
        if data:
            results.append(data)
    
    return results


def generate_dynamodb_items(parsed_data: list[dict]) -> dict:
    """Generate DynamoDB-ready items from parsed data."""
    curriculum_items = []
    topic_items = []
    subtopic_items = []
    
    for data in parsed_data:
        subject = data["subject"]
        year = data["year"]
        
        # Create curriculum entry for each grade
        for grade in data["grades"]:
            curriculum_id = f"{subject.lower().replace(' ', '-')}-grade{grade}-{year}"
            curriculum_id = re.sub(r'[^a-z0-9-]', '', curriculum_id)
            
            curriculum_items.append({
                "curriculumId": curriculum_id,
                "subjectName": subject,
                "grade": str(grade),
                "year": year,
                "description": f"{subject} curriculum for Grade {grade} ({year})",
            })
            
            # Create topics for this curriculum
            for i, topic in enumerate(data["topics"]):
                topic_id = generate_id("topic")
                
                topic_items.append({
                    "topicId": topic_id,
                    "curriculumId": curriculum_id,
                    "topicName": topic["topicName"],
                    "term": topic["term"] if topic["term"] > 0 else 1,
                    "orderIndex": i,
                    "weekCount": 2,  # Default, can be updated
                    "context": topic.get("context", ""),
                })
    
    return {
        "curriculum": curriculum_items,
        "topics": topic_items,
        "subtopics": subtopic_items,
    }


def main():
    """Main entry point."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    downloads_dir = os.path.join(script_dir, "Downloads")
    output_file = os.path.join(script_dir, "curriculum_data.json")
    
    print("="*60)
    print("ATP Parser - Extracting CAPS Curriculum Data")
    print("="*60)
    
    # Parse all PPTX files
    parsed_data = process_all_atps(downloads_dir)
    print(f"\nSuccessfully parsed {len(parsed_data)} files")
    
    # Generate DynamoDB items
    db_items = generate_dynamodb_items(parsed_data)
    
    print(f"\nGenerated items:")
    print(f"  - Curriculum entries: {len(db_items['curriculum'])}")
    print(f"  - Topic entries: {len(db_items['topics'])}")
    print(f"  - Subtopic entries: {len(db_items['subtopics'])}")
    
    # Save to JSON
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(db_items, f, indent=2, ensure_ascii=False)
    
    print(f"\nData saved to: {output_file}")
    
    # Also save parsed raw data for reference
    raw_output = os.path.join(script_dir, "atp_parsed_raw.json")
    with open(raw_output, 'w', encoding='utf-8') as f:
        json.dump(parsed_data, f, indent=2, ensure_ascii=False)
    
    print(f"Raw parsed data saved to: {raw_output}")


if __name__ == "__main__":
    main()
